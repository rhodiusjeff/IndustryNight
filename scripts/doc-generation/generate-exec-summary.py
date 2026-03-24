#!/usr/bin/env python3
"""
Industry Night — Executive Summary (Non-Technical)

Short 5-slide deck for stakeholder/investor-level audiences.
No code metrics, no architecture — just what we built, why, and where we're going.

Usage:
    source /tmp/pptx-env/bin/activate
    python3 scripts/doc-generation/generate-exec-summary.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import date

# --- Brand Colors ---------------------------------------------------------------
ACCENT      = RGBColor(0x6C, 0x5C, 0xE7)
ACCENT_LIGHT= RGBColor(0xA2, 0x96, 0xF0)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY    = RGBColor(0x88, 0x88, 0x99)
GREEN       = RGBColor(0x00, 0xB8, 0x94)
AMBER       = RGBColor(0xFD, 0xCB, 0x6E)
RED_SOFT    = RGBColor(0xE1, 0x7A, 0x7A)
CARD_BG     = RGBColor(0x22, 0x22, 0x3A)
SLIDE_BG    = RGBColor(0x12, 0x12, 0x22)

REPORT_DATE = date.today().strftime("%B %d, %Y")

# --- Helpers --------------------------------------------------------------------

def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def set_text(tf, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = alignment

def add_para(tf, text, size=14, color=WHITE, bold=False, space_before=Pt(4), alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.space_before = space_before
    p.space_after = Pt(2)
    p.alignment = alignment

def add_bullet(tf, text, size=14, color=WHITE, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.level = 0
    p.space_before = Pt(6)
    p.space_after = Pt(3)

def add_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def accent_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

def slide_header(slide, title, subtitle=None):
    set_slide_bg(slide, SLIDE_BG)
    accent_bar(slide)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.55), Inches(1.2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    tb = add_textbox(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7))
    set_text(tb.text_frame, title, size=32, color=WHITE, bold=True)
    if subtitle:
        tb2 = add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.4))
        set_text(tb2.text_frame, subtitle, size=15, color=MID_GRAY)

def add_stat_card(slide, left, top, number, label, accent_color=ACCENT):
    card = add_rounded_rect(slide, left, top, Inches(2.6), Inches(1.2))
    tf = card.text_frame
    tf.margin_top = Inches(0.15)
    tf.margin_left = Inches(0.2)
    tf.word_wrap = True
    set_text(tf, number, size=32, color=accent_color, bold=True)
    add_para(tf, label, size=11, color=LIGHT_GRAY, space_before=Pt(0))


# --- Presentation ---------------------------------------------------------------

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# === SLIDE 1: Title =============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, SLIDE_BG)
accent_bar(sl)

tb = add_textbox(sl, Inches(1.2), Inches(1.6), Inches(10), Inches(1.5))
set_text(tb.text_frame, "INDUSTRY NIGHT", size=56, color=WHITE, bold=True)
add_para(tb.text_frame, "The Home for Creative Professionals", size=24, color=ACCENT_LIGHT, space_before=Pt(12))

line = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(3.8), Inches(4), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

tb2 = add_textbox(sl, Inches(1.2), Inches(4.2), Inches(10), Inches(1.5))
set_text(tb2.text_frame, f"Executive Summary  |  {REPORT_DATE}", size=16, color=MID_GRAY)
add_para(tb2.text_frame, "Event-first social network for creative professionals", size=14, color=MID_GRAY, space_before=Pt(8))


# === SLIDE 2: The Opportunity ===================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(sl, "The Opportunity")

# Left: the problem
card1 = add_rounded_rect(sl, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.5))
tf = card1.text_frame
tf.margin_top = Inches(0.3)
tf.margin_left = Inches(0.35)
tf.word_wrap = True
set_text(tf, "What's Missing Today", size=22, color=ACCENT_LIGHT, bold=True)
add_para(tf, "", size=8, color=WHITE)
add_bullet(tf, "Creative professionals (stylists, photographers, makeup artists, videographers) have no dedicated networking platform", size=16, color=LIGHT_GRAY)
add_bullet(tf, "Existing tools are generic event sites or portfolio showcases -- none tie real-world attendance to community membership", size=16, color=LIGHT_GRAY)
add_bullet(tf, "Industry Night already runs successful live events with an established audience and social following", size=16, color=LIGHT_GRAY)

# Right: the solution
card2 = add_rounded_rect(sl, Inches(6.8), Inches(2.0), Inches(5.6), Inches(4.5))
tf2 = card2.text_frame
tf2.margin_top = Inches(0.3)
tf2.margin_left = Inches(0.35)
tf2.word_wrap = True
set_text(tf2, "What We're Building", size=22, color=ACCENT_LIGHT, bold=True)
add_para(tf2, "", size=8, color=WHITE)
add_bullet(tf2, "A mobile app and admin platform purpose-built for creative professionals", size=16, color=LIGHT_GRAY)
add_bullet(tf2, "Verified community: earn membership by attending events and making real QR connections", size=16, color=LIGHT_GRAY)
add_bullet(tf2, "Every connection proves physical co-presence -- impossible to fake, impossible for competitors to copy", size=16, color=LIGHT_GRAY)
add_bullet(tf2, "Sponsor perks and community feed keep members engaged between events", size=16, color=LIGHT_GRAY)


# === SLIDE 3: What We've Built ==================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(sl, "What We've Built", "From concept to working platform in 18 days")

# Three column layout
cols = [
    ("Mobile App", "Working", GREEN, [
        "Phone-based login (no passwords needed)",
        "Profile with specialties and bio",
        "Event browsing with ticket status",
        "Check-in with activation codes",
        "Instant QR connections with celebration",
        "Real-time connection notifications",
        "Remember-me login (stay signed in)",
    ]),
    ("Admin Dashboard", "Working", GREEN, [
        "Secure admin login (separate system)",
        "Full user management (add, view, moderate)",
        "Event creation with image uploads",
        "Ticket management (issue, refund, view)",
        "Sponsor and vendor management",
        "Activation code system",
        "Event publish gate (images + venue + Posh required)",
    ]),
    ("Platform", "Production-Ready", GREEN, [
        "Cloud infrastructure on AWS (auto-scaling)",
        "Posh.vip ticket integration (webhooks)",
        "SMS verification via Twilio",
        "S3 image storage with CDN",
        "Cost hibernation ($160/mo to $3/mo when idle)",
        "Database with full audit trail",
        "Comprehensive operational tooling",
    ]),
]

for i, (title, status, status_color, bullets) in enumerate(cols):
    left = Inches(0.8 + i * 4.1)
    card = add_rounded_rect(sl, left, Inches(2.0), Inches(3.8), Inches(5.0))
    tf = card.text_frame
    tf.margin_top = Inches(0.25)
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.2)
    tf.word_wrap = True
    set_text(tf, title, size=22, color=WHITE, bold=True)

    # Status badge
    badge = add_rounded_rect(sl, left + Inches(2.0), Inches(2.15), Inches(1.6), Inches(0.35), fill_color=status_color)
    btf = badge.text_frame
    btf.margin_top = Inches(0.0)
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(btf, status, size=10, color=SLIDE_BG, bold=True, alignment=PP_ALIGN.CENTER)

    for b in bullets:
        add_bullet(tf, b, size=13, color=LIGHT_GRAY)


# === SLIDE 4: Social Network & Sponsor Value ====================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(sl, "The Data Advantage", "Every connection proves two professionals physically met")

# Left: network value
card1 = add_rounded_rect(sl, Inches(0.8), Inches(2.0), Inches(5.6), Inches(2.6))
tf = card1.text_frame
tf.margin_top = Inches(0.25)
tf.margin_left = Inches(0.3)
tf.word_wrap = True
set_text(tf, "What Makes IN Different", size=18, color=ACCENT_LIGHT, bold=True)
add_bullet(tf, "QR connections = verified proof of physical meeting (no other platform has this)", size=13, color=LIGHT_GRAY)
add_bullet(tf, "We know each user's specialty, events attended, and who they connect with", size=13, color=LIGHT_GRAY)
add_bullet(tf, "Cross-specialty affinity data tells sponsors exactly where to focus", size=13, color=LIGHT_GRAY)
add_bullet(tf, "Verification ladder drives event attendance (attend + connect = unlock features)", size=13, color=LIGHT_GRAY)

# Right: sponsor tiers
card2 = add_rounded_rect(sl, Inches(6.8), Inches(2.0), Inches(5.6), Inches(2.6))
tf2 = card2.text_frame
tf2.margin_top = Inches(0.25)
tf2.margin_left = Inches(0.3)
tf2.word_wrap = True
set_text(tf2, "Sponsor Revenue Tiers", size=18, color=ACCENT_LIGHT, bold=True)
add_para(tf2, "", size=4, color=WHITE)

# Tier cards
for tier_label, price, desc, color in [
    ("Tier 1", "$500-2K/event", "Logo placement on event pages", MID_GRAY),
    ("Tier 2", "$2-5K/event", "Verified audience access + redemption tracking", AMBER),
    ("Tier 3", "$5-20K/quarter", "Ongoing audience intelligence & data partnership", GREEN),
]:
    add_para(tf2, f"{tier_label}:  {price}", size=13, color=color, bold=True, space_before=Pt(6))
    add_para(tf2, f"    {desc}", size=11, color=LIGHT_GRAY, space_before=Pt(0))

# Bottom: stat cards
add_stat_card(sl, Inches(0.8), Inches(5.0), "Verified", "Every connection = proof of meeting", accent_color=GREEN)
add_stat_card(sl, Inches(3.8), Inches(5.0), "Specialty Data", "Demographics by creative discipline", accent_color=ACCENT)
add_stat_card(sl, Inches(6.8), Inches(5.0), "Affinity Maps", "Who connects with whom and why", accent_color=ACCENT_LIGHT)
add_stat_card(sl, Inches(9.8), Inches(5.0), "Recurring $", "Quarterly data partnerships", accent_color=AMBER)


# === SLIDE 5: Where We're Going =================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(sl, "Where We're Going")

# Roadmap items
roadmap = [
    ("Now", "Community & Engagement", ACCENT,
     "Wire the community feed (backend is built), add push notifications, "
     "implement verification gating. Make the app worth opening between events."),
    ("Next", "Professional Utility", GREEN,
     "Connection-only messaging (DMs with people you've met), profile portfolios for work showcase, "
     "creative search by specialty. Build the daily-use professional tool."),
    ("Then", "Sponsor Revenue Engine", ACCENT_LIGHT,
     "Discount redemption tracking, analytics computation for audience intelligence, "
     "sponsor post-event reports. Prove ROI to unlock Tier 2-3 pricing."),
    ("Future", "Growth & Scale", MID_GRAY,
     "Market area expansion beyond NYC, mutual connections display, structured collaboration board, "
     "external sharing for viral growth. Network effects compound."),
]

for i, (phase, title, color, desc) in enumerate(roadmap):
    top = Inches(2.0 + i * 1.25)

    # Phase label
    label = add_rounded_rect(sl, Inches(0.8), top + Inches(0.05), Inches(1.0), Inches(0.7), fill_color=color)
    ltf = label.text_frame
    ltf.margin_top = Inches(0.0)
    ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(ltf, phase, size=14, color=SLIDE_BG, bold=True, alignment=PP_ALIGN.CENTER)

    # Content card
    card = add_rounded_rect(sl, Inches(2.0), top, Inches(10.5), Inches(0.9))
    tf = card.text_frame
    tf.margin_top = Inches(0.12)
    tf.margin_left = Inches(0.3)
    tf.word_wrap = True
    set_text(tf, title, size=18, color=WHITE, bold=True)
    add_para(tf, desc, size=14, color=LIGHT_GRAY, space_before=Pt(4))

    # Connector
    if i < 3:
        conn = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(1.25), top + Inches(0.8), Pt(3), Inches(0.5))
        conn.fill.solid()
        conn.fill.fore_color.rgb = ACCENT
        conn.line.fill.background()

# Bottom note
tb = add_textbox(sl, Inches(0.8), Inches(6.6), Inches(11), Inches(0.5))
set_text(tb.text_frame, "Platform built from scratch in 18 days. Foundation complete -- now focused on engagement and monetization.", size=14, color=MID_GRAY)


# --- Save -----------------------------------------------------------------------
output_path = "docs/executive/Industry Night - Executive Summary.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
