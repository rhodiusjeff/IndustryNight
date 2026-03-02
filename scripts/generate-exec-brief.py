#!/usr/bin/env python3
"""
Industry Night -- Executive Brief PowerPoint Generator

Generates a professional .pptx presentation from project data.
Re-run this script to regenerate the deck with updated metrics.

Usage:
    source /tmp/pptx-env/bin/activate
    python3 scripts/generate-exec-brief.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import date

# --- Brand Colors ---------------------------------------------------------------
BLACK       = RGBColor(0x0F, 0x0F, 0x0F)
DARK_BG     = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT      = RGBColor(0x6C, 0x5C, 0xE7)   # Purple accent
ACCENT_LIGHT= RGBColor(0xA2, 0x96, 0xF0)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY    = RGBColor(0x88, 0x88, 0x99)
GREEN       = RGBColor(0x00, 0xB8, 0x94)
AMBER       = RGBColor(0xFD, 0xCB, 0x6E)
RED_SOFT    = RGBColor(0xE1, 0x7A, 0x7A)
CARD_BG     = RGBColor(0x22, 0x22, 0x3A)
SLIDE_BG    = RGBColor(0x12, 0x12, 0x22)

# --- Report Data (update these for weekly refresh) ------------------------------
REPORT_DATE = date.today().strftime("%B %d, %Y")
PERIOD      = "Project Inception through Week 3"

# --- Helpers --------------------------------------------------------------------

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def set_text(tf, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return p

def add_para(tf, text, size=14, color=WHITE, bold=False, space_before=Pt(4), space_after=Pt(2), alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.space_after = space_after
    p.alignment = alignment
    return p

def add_bullet(tf, text, size=14, color=WHITE, level=0, bold=False, font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.level = level
    p.space_before = Pt(3)
    p.space_after = Pt(2)
    return p

def add_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_stat_card(slide, left, top, number, label, accent_color=ACCENT):
    card = add_rounded_rect(slide, left, top, Inches(2.1), Inches(1.2))
    tf = card.text_frame
    tf.margin_top = Inches(0.15)
    tf.margin_left = Inches(0.15)
    tf.word_wrap = True
    set_text(tf, number, size=32, color=accent_color, bold=True)
    add_para(tf, label, size=11, color=LIGHT_GRAY, space_before=Pt(0))

def add_accent_line(slide, top):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), top, Inches(1.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

def slide_title(slide, title, subtitle=None):
    set_slide_bg(slide, SLIDE_BG)
    add_accent_line(slide, Inches(0.5))
    tb = add_textbox(slide, Inches(0.6), Inches(0.6), Inches(8.5), Inches(0.6))
    set_text(tb.text_frame, title, size=28, color=WHITE, bold=True)
    if subtitle:
        tb2 = add_textbox(slide, Inches(0.6), Inches(1.15), Inches(8.5), Inches(0.35))
        set_text(tb2.text_frame, subtitle, size=13, color=MID_GRAY)


def add_table_slide(slide, title, headers, rows, col_widths=None, subtitle=None):
    """Add a styled table to a slide."""
    slide_title(slide, title, subtitle)

    num_rows = len(rows) + 1
    num_cols = len(headers)
    top = Inches(1.65) if subtitle else Inches(1.45)
    table_shape = slide.shapes.add_table(num_rows, num_cols, Inches(0.6), top, Inches(8.8), Inches(0.35 * num_rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = "Calibri"

    # Data rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG if r_idx % 2 == 0 else SLIDE_BG
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = WHITE
                p.font.name = "Calibri"


# --- Presentation ---------------------------------------------------------------

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# === SLIDE 1: Title =============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(sl, SLIDE_BG)

# Big accent bar
bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()

# Title
tb = add_textbox(sl, Inches(1.0), Inches(1.8), Inches(10), Inches(1.2))
set_text(tb.text_frame, "INDUSTRY NIGHT", size=52, color=WHITE, bold=True)
add_para(tb.text_frame, "Executive Brief", size=28, color=ACCENT_LIGHT, bold=False, space_before=Pt(8))

# Divider line
line = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.7), Inches(4), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

# Subtitle info
tb2 = add_textbox(sl, Inches(1.0), Inches(4.0), Inches(10), Inches(1.5))
set_text(tb2.text_frame, REPORT_DATE, size=16, color=LIGHT_GRAY)
add_para(tb2.text_frame, PERIOD, size=14, color=MID_GRAY, space_before=Pt(8))
add_para(tb2.text_frame, "Confidential", size=12, color=MID_GRAY, space_before=Pt(20))


# === SLIDE 2: What Is Industry Night? ===========================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(sl, "What Is Industry Night?", "Event-first social network for creative professionals")

# Left column -- description
tb = add_textbox(sl, Inches(0.6), Inches(1.6), Inches(5.5), Inches(4.5))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "The Problem", size=18, color=ACCENT_LIGHT, bold=True)
add_bullet(tf, "Creative professionals (stylists, makeup artists, photographers, videographers) lack a purpose-built networking platform", size=13, color=LIGHT_GRAY)
add_bullet(tf, "Existing tools are generic (Meetup, Eventbrite) or portfolio-focused (Behance)", size=13, color=LIGHT_GRAY)
add_bullet(tf, "No platform ties event attendance to verified community membership", size=13, color=LIGHT_GRAY)
add_para(tf, "", size=6, color=WHITE)
add_para(tf, "The Solution", size=18, color=ACCENT_LIGHT, bold=True)
add_bullet(tf, "Open registration with verification earned through real-world actions", size=13, color=LIGHT_GRAY)
add_bullet(tf, "QR-code networking at events creates instant, mutual connections", size=13, color=LIGHT_GRAY)
add_bullet(tf, "Verification through attendance builds trust and authenticity", size=13, color=LIGHT_GRAY)
add_bullet(tf, "Two apps: Social (for creatives) + Admin (for operators)", size=13, color=LIGHT_GRAY)

# Right column -- two product cards
card1 = add_rounded_rect(sl, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.2))
tf1 = card1.text_frame
tf1.margin_top = Inches(0.2)
tf1.margin_left = Inches(0.25)
tf1.word_wrap = True
set_text(tf1, "Social App", size=18, color=ACCENT_LIGHT, bold=True)
add_para(tf1, "Target:  Creative professionals", size=12, color=LIGHT_GRAY)
add_para(tf1, "Platforms:  iOS, Android, Web", size=12, color=LIGHT_GRAY)
add_para(tf1, "Purpose:  Attend events, QR networking, community feed, perks", size=12, color=LIGHT_GRAY)

card2 = add_rounded_rect(sl, Inches(6.8), Inches(4.1), Inches(5.8), Inches(2.2))
tf2 = card2.text_frame
tf2.margin_top = Inches(0.2)
tf2.margin_left = Inches(0.25)
tf2.word_wrap = True
set_text(tf2, "Admin App", size=18, color=ACCENT_LIGHT, bold=True)
add_para(tf2, "Target:  Platform operators", size=12, color=LIGHT_GRAY)
add_para(tf2, "Platforms:  Web, iOS, Android", size=12, color=LIGHT_GRAY)
add_para(tf2, "Purpose:  Manage events, users, tickets, sponsors, moderation", size=12, color=LIGHT_GRAY)


# === SLIDE 3: How It Works =====================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(sl, "How It Works", "Open registration to verified community member")

steps = [
    ("1", "Register", "Download app\nSMS phone login\nNo password needed", "Set up profile\nwith specialties"),
    ("2", "Get Ticket", "Purchase on Posh.vip\nor admin-issued", "Auto-linked to\naccount by phone"),
    ("3", "Attend Event", "Door staff gives\n4-digit activation code", "Check in via\napp (code + ticket)"),
    ("4", "Connect", "Scan QR codes for\ninstant connections", "First connection =\nVerified status"),
    ("5", "Full Access", "Community board,\nsponsor perks", "Networking persists\nacross events"),
]

for i, (num, title, line1, line2) in enumerate(steps):
    left = Inches(0.6 + i * 2.5)
    # Step number circle
    circle = sl.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.75), Inches(1.7), Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    ctf = circle.text_frame
    ctf.margin_top = Inches(0.0)
    ctf.margin_bottom = Inches(0.0)
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(ctf, num, size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Connector line (except last)
    if i < 4:
        conn = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(1.35), Inches(1.95), Inches(1.9), Pt(2))
        conn.fill.solid()
        conn.fill.fore_color.rgb = ACCENT
        conn.line.fill.background()

    # Step title
    tb = add_textbox(sl, left, Inches(2.4), Inches(2.2), Inches(0.4))
    set_text(tb.text_frame, title, size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Card
    card = add_rounded_rect(sl, left, Inches(2.9), Inches(2.2), Inches(1.8))
    ctf = card.text_frame
    ctf.margin_top = Inches(0.2)
    ctf.margin_left = Inches(0.15)
    ctf.word_wrap = True
    set_text(ctf, line1, size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_para(ctf, "", size=6, color=WHITE)
    add_para(ctf, line2, size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# === SLIDE 4: Project Timeline ==================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(sl, "Project Timeline", "18 days from first commit to full-stack working platform")

milestones = [
    ("Feb 4", "Requirements\nFinalized", "Product requirements\ndocument v1.8"),
    ("Feb 13", "Initial\nCommit", "Project scaffolded,\nmonorepo structure"),
    ("Feb 23", "Infrastructure\nLaydown", "AWS EKS, domain\nmigration, ops tooling"),
    ("Feb 25", "Foundation\nComplete", "PR #1 merged,\nadmin auth working"),
    ("Mar 1", "Full Stack\nWorking", "QR connections,\ntickets, check-in,\nadversarial review"),
]

# Timeline line
line = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.6), Inches(11.0), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

for i, (dt, title, desc) in enumerate(milestones):
    cx = Inches(1.0 + i * 2.6)
    # Dot
    dot = sl.shapes.add_shape(MSO_SHAPE.OVAL, cx + Inches(0.5), Inches(3.45), Inches(0.35), Inches(0.35))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT
    dot.line.fill.background()
    # Date above
    tb = add_textbox(sl, cx - Inches(0.2), Inches(2.3), Inches(1.8), Inches(0.35))
    set_text(tb.text_frame, dt, size=14, color=ACCENT_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    tb = add_textbox(sl, cx - Inches(0.2), Inches(2.65), Inches(1.8), Inches(0.7))
    set_text(tb.text_frame, title, size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Description card below
    card = add_rounded_rect(sl, cx - Inches(0.3), Inches(4.1), Inches(2.0), Inches(1.6))
    ctf = card.text_frame
    ctf.margin_top = Inches(0.15)
    ctf.margin_left = Inches(0.15)
    ctf.word_wrap = True
    set_text(ctf, desc, size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# === SLIDE 5: Codebase Metrics ==================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(sl, "Codebase at a Glance")

# Top row: 5 stat cards
stats = [
    ("~17,000", "Lines of Code"),
    ("120+", "Source Files"),
    ("22", "Database Tables"),
    ("35+", "App Screens"),
    ("11", "API Route Modules"),
]
for i, (num, label) in enumerate(stats):
    add_stat_card(sl, Inches(0.6 + i * 2.5), Inches(1.6), num, label)

# Bottom row: 4 stat cards
stats2 = [
    ("9", "Data Models"),
    ("7", "API Clients"),
    ("20+", "Git Commits"),
    ("15+", "Operational Scripts"),
]
for i, (num, label) in enumerate(stats2):
    add_stat_card(sl, Inches(0.6 + i * 2.5), Inches(3.2), num, label, accent_color=GREEN)

# LOC breakdown table
add_table_data = [
    ("Backend API (TypeScript)", "~2,500", "25+"),
    ("Flutter Apps + Shared (Dart)", "~11,000", "70+"),
    ("Database (SQL)", "700+", "6"),
    ("Scripts (JS + Bash + Python)", "3,000+", "20+"),
    ("Infrastructure (YAML)", "250+", "7"),
]

tbl_shape = sl.shapes.add_table(len(add_table_data)+1, 3, Inches(0.6), Inches(4.9), Inches(6.5), Inches(0.32 * (len(add_table_data)+1)))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(3.5)
tbl.columns[1].width = Inches(1.5)
tbl.columns[2].width = Inches(1.5)
headers = ["Component", "Lines of Code", "Files"]
for i, h in enumerate(headers):
    c = tbl.cell(0, i)
    c.text = h
    c.fill.solid()
    c.fill.fore_color.rgb = ACCENT
    for p in c.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = "Calibri"
for r, row in enumerate(add_table_data):
    for c, val in enumerate(row):
        cell = tbl.cell(r+1, c)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG if r % 2 == 0 else SLIDE_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"


# === SLIDE 6: What Has Been Built -- Backend ====================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(sl, "What Has Been Built", "Backend API + Database")

# Backend API card
card1 = add_rounded_rect(sl, Inches(0.6), Inches(1.6), Inches(6.0), Inches(5.2))
tf = card1.text_frame
tf.margin_top = Inches(0.2)
tf.margin_left = Inches(0.25)
tf.word_wrap = True
set_text(tf, "Backend API  (Node.js / Express / TypeScript)", size=16, color=ACCENT_LIGHT, bold=True)
add_bullet(tf, "11 route modules: auth, users, events, connections, posts, sponsors, vendors, discounts, webhooks, admin, admin-auth", size=12, color=LIGHT_GRAY)
add_bullet(tf, "4 middleware layers: JWT auth, admin auth (token family separation), role-based access, Zod validation", size=12, color=LIGHT_GRAY)
add_bullet(tf, "3 service integrations: Twilio SMS/Verify, AWS SES email, Posh.vip webhooks", size=12, color=LIGHT_GRAY)
add_bullet(tf, "JWT dual-auth: separate social and admin token families prevent cross-app token reuse", size=12, color=LIGHT_GRAY)
add_bullet(tf, "S3 image upload with public CDN URLs (event images, future: profile photos)", size=12, color=LIGHT_GRAY)
add_bullet(tf, "Auto-refresh tokens (15 min access + long-lived refresh)", size=12, color=LIGHT_GRAY)

# Database card
card2 = add_rounded_rect(sl, Inches(7.0), Inches(1.6), Inches(5.7), Inches(5.2))
tf2 = card2.text_frame
tf2.margin_top = Inches(0.2)
tf2.margin_left = Inches(0.25)
tf2.word_wrap = True
set_text(tf2, "Database  (PostgreSQL 15 on RDS)", size=16, color=ACCENT_LIGHT, bold=True)
add_bullet(tf2, "22 tables across 4 migrations", size=12, color=LIGHT_GRAY)
add_bullet(tf2, "Core: users, admin_users, events, tickets, posh_orders", size=12, color=LIGHT_GRAY)
add_bullet(tf2, "Social: connections, posts, comments, likes", size=12, color=LIGHT_GRAY)
add_bullet(tf2, "Business: sponsors, vendors, discounts, event_sponsors", size=12, color=LIGHT_GRAY)
add_bullet(tf2, "Media: event_images (up to 5 per event, hero image system)", size=12, color=LIGHT_GRAY)
add_bullet(tf2, "Operations: audit log, analytics tables, GDPR export tracking", size=12, color=LIGHT_GRAY)
add_bullet(tf2, "CASCADE delete design with audit log preservation", size=12, color=LIGHT_GRAY)


# === SLIDE 7: What Has Been Built -- Apps =======================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(sl, "What Has Been Built", "Social App + Admin App + Shared Package")

# Social App
card1 = add_rounded_rect(sl, Inches(0.6), Inches(1.6), Inches(4.0), Inches(5.2))
tf = card1.text_frame
tf.margin_top = Inches(0.2)
tf.margin_left = Inches(0.25)
tf.word_wrap = True
set_text(tf, "Social App  (Flutter/Dart)", size=15, color=ACCENT_LIGHT, bold=True)
add_para(tf, "19+ screens, 8 feature modules", size=11, color=MID_GRAY)
add_bullet(tf, "Auth: phone entry + SMS verify + remember-me", size=12, color=LIGHT_GRAY)
add_bullet(tf, "Events: browse, detail, tickets, check-in", size=12, color=LIGHT_GRAY)
add_bullet(tf, "Networking: QR display, scanner, instant connect", size=12, color=LIGHT_GRAY)
add_bullet(tf, "Real-time connection notifications (polling)", size=12, color=LIGHT_GRAY)
add_bullet(tf, "Celebration overlay on new connections", size=12, color=LIGHT_GRAY)
add_bullet(tf, "Profile: view, edit, settings", size=12, color=LIGHT_GRAY)
add_bullet(tf, "Community feed, search, perks (UI built, wiring pending)", size=12, color=LIGHT_GRAY)

# Admin App
card2 = add_rounded_rect(sl, Inches(4.9), Inches(1.6), Inches(4.0), Inches(5.2))
tf2 = card2.text_frame
tf2.margin_top = Inches(0.2)
tf2.margin_left = Inches(0.25)
tf2.word_wrap = True
set_text(tf2, "Admin App  (Flutter/Dart)", size=15, color=ACCENT_LIGHT, bold=True)
add_para(tf2, "16+ screens, 7 feature modules", size=11, color=MID_GRAY)
add_bullet(tf2, "Auth: email/password admin login", size=12, color=LIGHT_GRAY)
add_bullet(tf2, "Dashboard: stats overview", size=12, color=LIGHT_GRAY)
add_bullet(tf2, "Users: list, detail, add, manage tickets", size=12, color=LIGHT_GRAY)
add_bullet(tf2, "Events: full lifecycle + image management", size=12, color=LIGHT_GRAY)
add_bullet(tf2, "Publish gate: Posh ID + venue + images required", size=12, color=LIGHT_GRAY)
add_bullet(tf2, "Sponsors: link to events, discount management", size=12, color=LIGHT_GRAY)
add_bullet(tf2, "Global image catalog with bulk operations", size=12, color=LIGHT_GRAY)

# Shared Package
card3 = add_rounded_rect(sl, Inches(9.2), Inches(1.6), Inches(3.5), Inches(5.2))
tf3 = card3.text_frame
tf3.margin_top = Inches(0.2)
tf3.margin_left = Inches(0.25)
tf3.word_wrap = True
set_text(tf3, "Shared Package", size=15, color=ACCENT_LIGHT, bold=True)
add_para(tf3, "Reused by both apps", size=11, color=MID_GRAY)
add_bullet(tf3, "9 data models with JSON serialization", size=12, color=LIGHT_GRAY)
add_bullet(tf3, "7 API clients (base HTTP + typed endpoints)", size=12, color=LIGHT_GRAY)
add_bullet(tf3, "Secure token storage", size=12, color=LIGHT_GRAY)
add_bullet(tf3, "Phone validators", size=12, color=LIGHT_GRAY)
add_bullet(tf3, "Display formatters", size=12, color=LIGHT_GRAY)
add_bullet(tf3, "Constants / enums", size=12, color=LIGHT_GRAY)


# === SLIDE 8: Infrastructure & Operations =======================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(sl, "Infrastructure & Operations", "AWS EKS + Kubernetes + operational tooling")

# Infra card
card1 = add_rounded_rect(sl, Inches(0.6), Inches(1.6), Inches(5.8), Inches(3.0))
tf = card1.text_frame
tf.margin_top = Inches(0.2)
tf.margin_left = Inches(0.25)
tf.word_wrap = True
set_text(tf, "Cloud Infrastructure  (AWS)", size=16, color=ACCENT_LIGHT, bold=True)
add_bullet(tf, "EKS cluster (us-east-1) with auto-scaling 2-10 pods", size=12, color=LIGHT_GRAY)
add_bullet(tf, "Kubernetes: deployment, service, ingress, secrets, HPA", size=12, color=LIGHT_GRAY)
add_bullet(tf, "SSL/TLS via ACM on api.industrynight.net", size=12, color=LIGHT_GRAY)
add_bullet(tf, "ECR container registry with Docker build pipeline", size=12, color=LIGHT_GRAY)
add_bullet(tf, "RDS PostgreSQL 15 (force SSL, automated backups)", size=12, color=LIGHT_GRAY)
add_bullet(tf, "S3 bucket for event images (public-read ACL)", size=12, color=LIGHT_GRAY)

# COOP card
card2 = add_rounded_rect(sl, Inches(6.8), Inches(1.6), Inches(5.8), Inches(3.0))
tf2 = card2.text_frame
tf2.margin_top = Inches(0.2)
tf2.margin_left = Inches(0.25)
tf2.word_wrap = True
set_text(tf2, "COOP System  (Cost Optimization)", size=16, color=ACCENT_LIGHT, bold=True)
add_bullet(tf2, "Full teardown: EKS + RDS removed, data preserved in S3", size=12, color=LIGHT_GRAY)
add_bullet(tf2, "Full rebuild: infra recreated from scratch + data restored", size=12, color=LIGHT_GRAY)
add_bullet(tf2, "Database backup/restore via pg_dump (full + per-table)", size=12, color=LIGHT_GRAY)
add_bullet(tf2, "Single entry point: scripts/coop/coop.sh", size=12, color=LIGHT_GRAY)

# Cost cards
add_stat_card(sl, Inches(0.6), Inches(5.0), "~$160/mo", "Full Running", accent_color=AMBER)
add_stat_card(sl, Inches(3.1), Inches(5.0), "~$3/mo", "Hibernated (COOP)", accent_color=GREEN)

# Scripts card
card3 = add_rounded_rect(sl, Inches(5.6), Inches(5.0), Inches(7.0), Inches(1.8))
tf3 = card3.text_frame
tf3.margin_top = Inches(0.15)
tf3.margin_left = Inches(0.25)
tf3.word_wrap = True
set_text(tf3, "Operational Scripts (15+)", size=14, color=ACCENT_LIGHT, bold=True)
add_bullet(tf3, "seed-admin.js, migrate.js, db-reset.js, db-scrub-user.js (GDPR)", size=11, color=LIGHT_GRAY)
add_bullet(tf3, "deploy-api.sh, maintenance.sh, setup-local.sh", size=11, color=LIGHT_GRAY)
add_bullet(tf3, "COOP: status, teardown, rebuild, export, import", size=11, color=LIGHT_GRAY)
add_bullet(tf3, "generate-exec-brief.py, generate-exec-summary.py", size=11, color=LIGHT_GRAY)


# === SLIDE 9: Architecture Decisions ============================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_table_slide(sl, "Architecture Decisions",
    ["Decision", "Choice", "Rationale"],
    [
        ["Registration model", "Open registration + verification ladder", "Events list is marketing funnel; trust earned through attendance + connections"],
        ["Auth (social)", "Phone + SMS OTP", "Passwordless, phone-based identity for creatives"],
        ["Auth (admin)", "Email + password", "Separate admin_users table, separate token family"],
        ["Mobile framework", "Flutter/Dart", "Single codebase for iOS, Android, Web"],
        ["Ticketing", "Posh.vip (webhooks) + admin-issued", "Existing brand presence; auto-link by phone on registration"],
        ["QR networking", "Instant mutual connection (no request/accept)", "Reduces friction; trust earned through physical co-presence"],
        ["Orchestration", "AWS EKS (Kubernetes)", "Scalability, infrastructure learning"],
        ["State mgmt", "Provider + ChangeNotifier", "Simple, sufficient for current scope"],
        ["Cost mgmt", "COOP teardown/rebuild", "Reduce AWS spend during downtime"],
    ],
    col_widths=[2.2, 2.8, 3.8]
)


# === SLIDE 10: Implementation Status ============================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(sl, "Implementation Status vs Plan")

phases = [
    ("1A", "Foundation (Backend + Auth)", "Complete", GREEN),
    ("1B", "Core Mobile App", "Complete", GREEN),
    ("1C", "Verification & QR Networking", "Complete", GREEN),
    ("1D", "Event Social Features", "Screens Built", AMBER),
    ("1E", "Community Board", "Backend Only", AMBER),
    ("1F", "Creative Search", "Backend Only", AMBER),
    ("2A", "Admin App - Foundation", "Complete", GREEN),
    ("2B", "Admin App - Event Mgmt", "Complete", GREEN),
    ("2C", "Admin App - Sponsor Mgmt", "Complete", GREEN),
    ("2D", "Admin App - Vendor Mgmt", "Complete", GREEN),
    ("2E", "Admin App - Moderation", "Screens Built", AMBER),
    ("--", "Adversarial Review & Decisions", "Complete", GREEN),
    ("3", "Community Feed + Engagement", "Planned", MID_GRAY),
]

# Render as two columns of cards
for i, (phase, desc, status, color) in enumerate(phases):
    col = 0 if i < 7 else 1
    row = i if i < 7 else i - 7
    left = Inches(0.6) if col == 0 else Inches(6.8)
    top = Inches(1.6 + row * 0.73)

    card = add_rounded_rect(sl, left, top, Inches(5.8), Inches(0.6))
    tf = card.text_frame
    tf.margin_top = Inches(0.08)
    tf.margin_left = Inches(0.15)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.name = "Calibri"

    run1 = p.add_run()
    run1.text = f"  {phase}  "
    run1.font.size = Pt(11)
    run1.font.color.rgb = MID_GRAY
    run1.font.bold = True
    run1.font.name = "Calibri"

    run2 = p.add_run()
    run2.text = f"  {desc}"
    run2.font.size = Pt(12)
    run2.font.color.rgb = WHITE
    run2.font.name = "Calibri"

    # Status badge
    badge = add_rounded_rect(sl, left + Inches(4.3), top + Inches(0.1), Inches(1.3), Inches(0.38), fill_color=color)
    btf = badge.text_frame
    btf.margin_top = Inches(0.0)
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(btf, status, size=9, color=SLIDE_BG, bold=True, alignment=PP_ALIGN.CENTER)


# === SLIDE 11: Key Decisions Made ===============================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_table_slide(sl, "Strategic Decisions (from Adversarial Review)",
    ["Question", "Decision", "Implication"],
    [
        ["Invite-only or open?", "Open registration", "Events list = marketing funnel; verification ladder = trust gate"],
        ["Posh webhook creates users?", "No -- auto-link by phone", "Users create own accounts; Posh orders reconciled on registration"],
        ["Verification feature gating?", "Yes -- backend required", "requireVerified middleware gates community board, perks"],
        ["Who's Going / Who's Here?", "Build behind feature flag", "Product owner decides visibility; avoids selective attendance concern"],
        ["Server-side connection validation?", "Defer (low risk now)", "Client-side gate sufficient at current scale"],
        ["Activation code time window?", "Event lifecycle IS the gate", "No code_valid_start/end columns needed"],
        ["Market area filtering?", "Add to events + users", "market_area enum; default to user's home market"],
        ["Sponsor/vendor/perks model?", "Deferred to product owner", "Current CRUD adequate for MVP; details TBD"],
    ],
    col_widths=[2.5, 2.0, 4.3],
    subtitle="8 strategic questions resolved during requirements-vs-reality audit"
)


# === SLIDE 12: Current WIP + What's Next =======================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(sl, "Current Work & What's Next")

# What's done recently
card1 = add_rounded_rect(sl, Inches(0.6), Inches(1.6), Inches(5.8), Inches(4.5))
tf = card1.text_frame
tf.margin_top = Inches(0.2)
tf.margin_left = Inches(0.25)
tf.word_wrap = True
set_text(tf, "Recently Completed", size=18, color=GREEN, bold=True)
add_para(tf, "", size=4, color=WHITE)
add_bullet(tf, "Instant QR connections with celebration overlay", size=13, color=LIGHT_GRAY)
add_bullet(tf, "Real-time connection notifications (polling)", size=13, color=LIGHT_GRAY)
add_bullet(tf, "Ticket-gated check-in system", size=13, color=LIGHT_GRAY)
add_bullet(tf, "Admin ticket management (issue, refund, delete)", size=13, color=LIGHT_GRAY)
add_bullet(tf, "JWT auto-refresh (15-min token lifecycle)", size=13, color=LIGHT_GRAY)
add_bullet(tf, "Remember-me login persistence", size=13, color=LIGHT_GRAY)
add_bullet(tf, "Adversarial review: 8 decisions, 30 issues tracked", size=13, color=LIGHT_GRAY)
add_bullet(tf, "Social network analysis & sponsor revenue strategy", size=13, color=LIGHT_GRAY)

# What's Next
card2 = add_rounded_rect(sl, Inches(6.8), Inches(1.6), Inches(5.8), Inches(4.5))
tf2 = card2.text_frame
tf2.margin_top = Inches(0.2)
tf2.margin_left = Inches(0.25)
tf2.word_wrap = True
set_text(tf2, "Up Next", size=18, color=ACCENT_LIGHT, bold=True)
add_para(tf2, "", size=4, color=WHITE)
items = [
    ("1.", "Wire community feed", "Backend exists; connect Flutter UI to real data"),
    ("2.", "Push notifications", "FCM/APNs for events, connections, feed activity"),
    ("3.", "Verification gating", "Backend middleware to gate community + perks"),
    ("4.", "Connection-only DMs", "Messaging between people who physically met"),
    ("5.", "Redemption tracking", "Prove sponsor ROI: track discount code usage"),
]
for num, title, desc in items:
    add_para(tf2, f"{num}  {title}", size=14, color=WHITE, bold=True, space_before=Pt(10))
    add_para(tf2, f"     {desc}", size=11, color=MID_GRAY, space_before=Pt(0))


# === SLIDE 13: Scorecard ========================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(sl, "MVP Scorecard", "33 requirements items: 13 done, 5 partial-adequate, 7 tracked, 2 retired")

# Stat cards
add_stat_card(sl, Inches(0.6), Inches(1.5), "55%", "Done or Adequate", accent_color=GREEN)
add_stat_card(sl, Inches(3.1), Inches(1.5), "85%", "Resolved or Tracked", accent_color=ACCENT)
add_stat_card(sl, Inches(5.6), Inches(1.5), "30", "GitHub Issues Tracked", accent_color=ACCENT_LIGHT)
add_stat_card(sl, Inches(8.1), Inches(1.5), "8", "Strategic Decisions Made", accent_color=AMBER)

# Key gaps
card = add_rounded_rect(sl, Inches(0.6), Inches(3.0), Inches(12.0), Inches(3.7))
tf = card.text_frame
tf.margin_top = Inches(0.2)
tf.margin_left = Inches(0.25)
tf.word_wrap = True
set_text(tf, "Key Remaining Gaps (tracked in GitHub)", size=16, color=ACCENT_LIGHT, bold=True)
add_para(tf, "", size=4, color=WHITE)

gaps = [
    ("#18", "Wire community feed to API", "P0 -- highest retention impact"),
    ("#14", "Verification-based feature gating", "P1 -- gates community + perks access"),
    ("#12-13", "Posh phone normalization + auto-link", "P1 -- enables walk-in ticket flow"),
    ("#20", "Pre-MVP security review", "Required before public launch"),
    ("New", "Push notifications", "P0 -- no way to pull users back between events"),
    ("New", "Discount redemption tracking", "P1 -- enables Tier 2 sponsor revenue"),
    ("New", "Connection-only DMs", "P1 -- biggest functional gap for retention"),
]
for issue, desc, priority in gaps:
    add_para(tf, f"  {issue}:  {desc}  --  {priority}", size=12, color=LIGHT_GRAY, space_before=Pt(4))


# === SLIDE 14: Summary ==========================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, SLIDE_BG)

bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()

tb = add_textbox(sl, Inches(1.0), Inches(1.2), Inches(11), Inches(1.0))
set_text(tb.text_frame, "Summary", size=36, color=WHITE, bold=True)

tb2 = add_textbox(sl, Inches(1.0), Inches(2.2), Inches(10), Inches(4.5))
tf = tb2.text_frame
tf.word_wrap = True
set_text(tf, "In 18 days, we have built:", size=20, color=LIGHT_GRAY)
add_para(tf, "", size=10, color=WHITE)
bullets = [
    "A complete backend API with 11 route modules, dual JWT auth, and 3 external integrations",
    "A 22-table PostgreSQL database with 4 migrations, full audit trail, and analytics schema",
    "A social app with QR-based instant connections, real-time notifications, and ticket-gated check-in",
    "An admin app with full event lifecycle, image management, ticket management, and publish gate",
    "A shared Dart package with 9 models, 7 API clients, and reusable utilities",
    "Production AWS infrastructure (EKS, RDS, ECR, S3, SSL) with cost-optimization tooling",
    "A comprehensive adversarial review with 8 strategic decisions and 30 tracked issues",
    "A social network analysis with 3-tier sponsor revenue strategy",
]
for b in bullets:
    add_bullet(tf, b, size=14, color=LIGHT_GRAY)

line = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(6.3), Inches(4), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

tb3 = add_textbox(sl, Inches(1.0), Inches(6.5), Inches(10), Inches(0.5))
set_text(tb3.text_frame, f"~17,000 lines of code  |  120+ files  |  {REPORT_DATE}", size=13, color=MID_GRAY)


# --- Save -----------------------------------------------------------------------
output_path = "docs/executive/Industry Night - Executive Brief.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
