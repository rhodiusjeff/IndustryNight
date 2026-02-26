#!/usr/bin/env python3
"""
Industry Night — Executive Summary (Non-Technical)

Short 4-slide deck for stakeholder/investor-level audiences.
No code metrics, no architecture — just what we built, why, and where we're going.

Usage:
    source /tmp/pptx-env/bin/activate
    python3 scripts/generate-exec-summary.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import date

# ─── Brand Colors ───────────────────────────────────────────────────────────
ACCENT      = RGBColor(0x6C, 0x5C, 0xE7)
ACCENT_LIGHT= RGBColor(0xA2, 0x96, 0xF0)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY    = RGBColor(0x88, 0x88, 0x99)
GREEN       = RGBColor(0x00, 0xB8, 0x94)
AMBER       = RGBColor(0xFD, 0xCB, 0x6E)
CARD_BG     = RGBColor(0x22, 0x22, 0x3A)
SLIDE_BG    = RGBColor(0x12, 0x12, 0x22)

REPORT_DATE = date.today().strftime("%B %d, %Y")

# ─── Helpers ────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def set_text(tf, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = alignment

def add_para(tf, text, size=14, color=WHITE, bold=False, space_before=Pt(4), alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.space_before = space_before
    p.space_after = Pt(2)
    p.alignment = alignment

def add_bullet(tf, text, size=14, color=WHITE, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.level = 0
    p.space_before = Pt(6)
    p.space_after = Pt(3)

def add_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def accent_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

def slide_header(slide, title, subtitle=None):
    set_slide_bg(slide, SLIDE_BG)
    accent_bar(slide)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.55), Inches(1.2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    tb = add_textbox(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7))
    set_text(tb.text_frame, title, size=32, color=WHITE, bold=True)
    if subtitle:
        tb2 = add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.4))
        set_text(tb2.text_frame, subtitle, size=15, color=MID_GRAY)


# ─── Presentation ───────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# ━━━ SLIDE 1: Title ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, SLIDE_BG)
accent_bar(sl)

tb = add_textbox(sl, Inches(1.2), Inches(1.6), Inches(10), Inches(1.5))
set_text(tb.text_frame, "INDUSTRY NIGHT", size=56, color=WHITE, bold=True)
add_para(tb.text_frame, "Building the Home for Creative Professionals", size=24, color=ACCENT_LIGHT, space_before=Pt(12))

line = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(3.8), Inches(4), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

tb2 = add_textbox(sl, Inches(1.2), Inches(4.2), Inches(10), Inches(1.5))
set_text(tb2.text_frame, f"Executive Summary  |  {REPORT_DATE}", size=16, color=MID_GRAY)


# ━━━ SLIDE 2: The Opportunity ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(sl, "The Opportunity")

# Left: the problem
card1 = add_rounded_rect(sl, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.5))
tf = card1.text_frame
tf.margin_top = Inches(0.3)
tf.margin_left = Inches(0.35)
tf.word_wrap = True
set_text(tf, "What's Missing Today", size=22, color=ACCENT_LIGHT, bold=True)
add_para(tf, "", size=8, color=WHITE)
add_bullet(tf, "Creative professionals in NYC (stylists, photographers, makeup artists, videographers) have no dedicated networking platform", size=16, color=LIGHT_GRAY)
add_bullet(tf, "Existing tools are generic event sites or portfolio showcases — none tie real-world attendance to community membership", size=16, color=LIGHT_GRAY)
add_bullet(tf, "Industry Night already runs successful live events with an established audience and Instagram following", size=16, color=LIGHT_GRAY)

# Right: the solution
card2 = add_rounded_rect(sl, Inches(6.8), Inches(2.0), Inches(5.6), Inches(4.5))
tf2 = card2.text_frame
tf2.margin_top = Inches(0.3)
tf2.margin_left = Inches(0.35)
tf2.word_wrap = True
set_text(tf2, "What We're Building", size=22, color=ACCENT_LIGHT, bold=True)
add_para(tf2, "", size=8, color=WHITE)
add_bullet(tf2, "An invite-only mobile app and admin platform purpose-built for creative professionals", size=16, color=LIGHT_GRAY)
add_bullet(tf2, "Verified community: you earn membership by attending events and making real connections", size=16, color=LIGHT_GRAY)
add_bullet(tf2, "QR networking at events creates instant, mutual professional connections", size=16, color=LIGHT_GRAY)
add_bullet(tf2, "Sponsor perks and community feed keep members engaged between events", size=16, color=LIGHT_GRAY)


# ━━━ SLIDE 3: What We've Accomplished ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(sl, "What We've Accomplished", "From concept to working platform in two weeks")

# Three column layout
cols = [
    ("Mobile App", "Ready for Creatives", GREEN, [
        "Phone-based login (no passwords)",
        "Profile creation with specialties and social links",
        "Event browsing and check-in with activation codes",
        "QR code networking at events",
        "Community feed for posts and job listings",
        "Sponsor perks and discount codes",
        "User discovery and search by specialty",
    ]),
    ("Admin Dashboard", "Ready for Operators", GREEN, [
        "Secure admin login (separate from user accounts)",
        "User management — view, add, and moderate",
        "Event creation with activation code system",
        "Sponsor and vendor management",
        "Content moderation and announcements",
        "Dashboard with platform overview",
    ]),
    ("Platform & Ops", "Production-Ready", AMBER, [
        "Secure backend API handling all business logic",
        "Cloud infrastructure on AWS with auto-scaling",
        "Posh.vip ticket integration (webhook-based)",
        "SMS verification via Twilio",
        "Cost management: infrastructure can hibernate when not needed (~$3/mo vs ~$160/mo)",
        "Database with full audit trail",
    ]),
]

for i, (title, status, status_color, bullets) in enumerate(cols):
    left = Inches(0.8 + i * 4.1)
    card = add_rounded_rect(sl, left, Inches(2.0), Inches(3.8), Inches(5.0))
    tf = card.text_frame
    tf.margin_top = Inches(0.25)
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.2)
    tf.word_wrap = True
    set_text(tf, title, size=22, color=WHITE, bold=True)

    # Status badge
    badge = add_rounded_rect(sl, left + Inches(2.0), Inches(2.15), Inches(1.6), Inches(0.35), fill_color=status_color)
    btf = badge.text_frame
    btf.margin_top = Inches(0.0)
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(btf, status, size=10, color=SLIDE_BG, bold=True, alignment=PP_ALIGN.CENTER)

    for b in bullets:
        add_bullet(tf, b, size=13, color=LIGHT_GRAY)


# ━━━ SLIDE 4: Where We're Going ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(sl, "Where We're Going")

# Roadmap items as horizontal cards
roadmap = [
    ("Now", "Launch Readiness", ACCENT,
     "Complete admin login, end-to-end testing, prepare for first live event trial with real users."),
    ("Next", "First Live Event", GREEN,
     "Test the full experience at an Industry Night: ticket purchase, app download, activation code, QR networking, verified status."),
    ("Then", "Growth Features", ACCENT_LIGHT,
     "In-app ticket purchasing (replace Posh), push notifications for event reminders, analytics dashboard for event performance."),
    ("Future", "Scale & Expand", MID_GRAY,
     "Expand beyond NYC. Deeper sponsor integrations. Enhanced community features. Native Android optimizations."),
]

for i, (phase, title, color, desc) in enumerate(roadmap):
    top = Inches(2.0 + i * 1.25)

    # Phase label
    label = add_rounded_rect(sl, Inches(0.8), top + Inches(0.05), Inches(1.0), Inches(0.7), fill_color=color)
    ltf = label.text_frame
    ltf.margin_top = Inches(0.0)
    ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(ltf, phase, size=14, color=SLIDE_BG, bold=True, alignment=PP_ALIGN.CENTER)

    # Content card
    card = add_rounded_rect(sl, Inches(2.0), top, Inches(10.5), Inches(0.9))
    tf = card.text_frame
    tf.margin_top = Inches(0.12)
    tf.margin_left = Inches(0.3)
    tf.word_wrap = True
    set_text(tf, title, size=18, color=WHITE, bold=True)
    add_para(tf, desc, size=14, color=LIGHT_GRAY, space_before=Pt(4))

    # Connector line between rows (except last)
    if i < 3:
        conn = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(1.25), top + Inches(0.8), Pt(3), Inches(0.5))
        conn.fill.solid()
        conn.fill.fore_color.rgb = ACCENT
        conn.line.fill.background()

# Bottom note
tb = add_textbox(sl, Inches(0.8), Inches(6.6), Inches(11), Inches(0.5))
set_text(tb.text_frame, "Platform was built from scratch in 12 days. Foundation is complete — now focused on launch readiness.", size=14, color=MID_GRAY)


# ─── Save ───────────────────────────────────────────────────────────────────
output_path = "docs/Industry Night - Executive Summary.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
